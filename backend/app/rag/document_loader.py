import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Set DEBUG = True to print chunk samples and detailed file info
DEBUG = True


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".png", ".jpg", ".jpeg"}


def load_documents(app_path: str) -> list:
    """
    Load ALL documents from an application folder.
    Supports: PDF, DOCX, XLSX/XLS, PPTX, PNG/JPG/JPEG (OCR)
    Returns list of {"text": str, "source": str}
    """
    docs = []
    path = Path(app_path)

    if not path.exists():
        logger.error(f"[LOADER] Path does not exist: {app_path}")
        return docs

    all_files = [f for f in path.rglob("*") if f.is_file()]
    supported = [f for f in all_files if f.suffix.lower() in SUPPORTED_EXTENSIONS]
    skipped   = [f.name for f in all_files if f.suffix.lower() not in SUPPORTED_EXTENSIONS and f.suffix != ""]

    logger.info(f"[LOADER] Scanning: {app_path}")
    logger.info(f"[LOADER] Total files found: {len(all_files)}")
    logger.info(f"[LOADER] Supported files: {[f.name for f in supported]}")
    if skipped:
        logger.info(f"[LOADER] Skipped (unsupported): {skipped}")

    for file in supported:
        ext = file.suffix.lower()
        logger.info(f"[LOADER] Processing: {file.name}")

        if ext == ".pdf":
            result = _load_pdf(file)
        elif ext == ".docx":
            result = _load_docx(file)
        elif ext in (".xlsx", ".xls"):
            result = _load_excel(file)
        elif ext == ".pptx":
            result = _load_pptx(file)
        elif ext in (".png", ".jpg", ".jpeg"):
            result = _load_image(file)
        else:
            result = []

        if result:
            docs.extend(result)
            logger.info(f"[LOADER] ✓ {file.name} → {len(result)} segment(s)")
        else:
            logger.warning(f"[LOADER] ✗ {file.name} → 0 segments (empty or failed)")

    logger.info(f"[LOADER] Total raw segments loaded: {len(docs)}")

    if DEBUG and docs:
        logger.debug(f"[LOADER] Sample segment (first 200 chars): {docs[0]['text'][:200]!r}")

    return docs


def _load_pdf(file: Path) -> list:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(file))
        docs = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                docs.append({"text": text, "source": f"{file.name} [page {i+1}]"})
        logger.info(f"[PDF] {file.name}: {len(reader.pages)} pages, {len(docs)} non-empty")
        return docs
    except Exception as e:
        logger.error(f"[PDF] Failed to load {file.name}: {e}", exc_info=True)
        return []


def _load_docx(file: Path) -> list:
    try:
        from docx import Document
        doc = Document(str(file))
        segments = []

        # Paragraphs
        para_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())

        # Tables
        table_rows = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_rows.append(row_text)
        table_text = "\n".join(table_rows)

        full_text = "\n\n".join(filter(None, [para_text, table_text])).strip()
        if full_text:
            segments.append({"text": full_text, "source": file.name})

        logger.info(f"[DOCX] {file.name}: {len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables")
        return segments
    except Exception as e:
        logger.error(f"[DOCX] Failed to load {file.name}: {e}", exc_info=True)
        return []


def _load_excel(file: Path) -> list:
    try:
        import pandas as pd
        xl = pd.ExcelFile(str(file))
        docs = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            df = df.dropna(how="all").fillna("")
            if df.empty:
                logger.info(f"[XLSX] {file.name} sheet '{sheet}': empty, skipping")
                continue

            # Convert each row to readable text: "Col1: val1 | Col2: val2"
            rows = []
            for _, row in df.iterrows():
                row_parts = [f"{col}: {str(val).strip()}" for col, val in row.items() if str(val).strip()]
                if row_parts:
                    rows.append(" | ".join(row_parts))

            text = f"[Sheet: {sheet}]\n" + "\n".join(rows)
            if text.strip():
                docs.append({"text": text, "source": f"{file.name} [{sheet}]"})
                logger.info(f"[XLSX] {file.name} sheet '{sheet}': {len(rows)} rows")

        return docs
    except Exception as e:
        logger.error(f"[XLSX] Failed to load {file.name}: {e}", exc_info=True)
        return []


def _load_pptx(file: Path) -> list:
    """
    Extract text from all slides in a PPTX file.
    Uses shape.text (not runs) to capture all text regardless of formatting.
    Produces one document segment per slide that has content.
    """
    try:
        from pptx import Presentation

        logger.info(f"[PPT] Processing PPT: {file.name}")
        prs = Presentation(str(file))
        docs = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []

            for shape in slide.shapes:
                # Use shape.text — catches all text regardless of run structure
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                raw = f"[Slide {slide_num}]\n" + "\n".join(slide_texts)
                text = _clean_text(raw)
                if text:
                    docs.append({
                        "text": text,
                        "source": f"{file.name} [slide {slide_num}]",
                    })

        logger.info(f"[PPT] {file.name}: {len(prs.slides)} total slides, {len(docs)} non-empty")

        if not docs:
            logger.warning(
                f"[PPT] {file.name}: 0 text segments extracted. "
                "The file may contain only images or embedded objects."
            )

        return docs

    except ImportError:
        logger.error(
            "[PPT] python-pptx is not installed. "
            "Run: pip install python-pptx==0.6.23"
        )
        return []
    except Exception as e:
        logger.error(f"[PPT] Failed to load {file.name}: {e}", exc_info=True)
        return []


def _load_image(file: Path) -> list:
    """Extract text from an image using OCR (pytesseract). Fails gracefully."""
    try:
        import pytesseract
        from PIL import Image

        # Allow overriding tesseract binary path via env var (for portable installs)
        import os
        tesseract_cmd = os.getenv("TESSERACT_CMD", "")
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

        logger.info(f"[OCR] Processing Image: {file.name}")
        img = Image.open(str(file))

        # Convert to RGB if needed (handles RGBA, palette modes, etc.)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        raw_text = pytesseract.image_to_string(img, timeout=30)
        text = _clean_text(raw_text)

        logger.info(f"[OCR] {file.name}: OCR text length: {len(text)} chars")

        if not text:
            logger.warning(f"[OCR] {file.name}: No text extracted — skipping")
            return []

        return [{"text": text, "source": f"{file.name} [ocr]"}]

    except ImportError:
        logger.warning(
            "[OCR] pytesseract or Pillow not installed — skipping image. "
            "Install with: pip install pytesseract pillow"
        )
        return []
    except RuntimeError as e:
        # pytesseract raises RuntimeError on timeout
        logger.warning(f"[OCR] {file.name}: OCR timed out — skipping")
        return []
    except Exception as e:
        logger.warning(f"[OCR] {file.name}: OCR failed ({e}) — skipping gracefully")
        return []


def _clean_text(text: str) -> str:
    """Normalize whitespace, remove empty lines and noise."""
    import re
    if not text:
        return ""
    # Collapse multiple spaces/tabs to single space
    text = re.sub(r"[ \t]+", " ", text)
    # Remove lines that are purely whitespace or single characters
    lines = [ln.strip() for ln in text.splitlines()]
    lines = [ln for ln in lines if len(ln) > 1]
    # Collapse 3+ consecutive blank lines to one
    result = re.sub(r"\n{3,}", "\n\n", "\n".join(lines))
    return result.strip()
